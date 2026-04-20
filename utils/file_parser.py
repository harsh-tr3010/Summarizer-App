import PyPDF2
import docx
from pptx import Presentation
from PIL import Image
import pytesseract
import tempfile
from pdf2image import convert_from_bytes


def extract_text(uploaded_file):
    file_type = uploaded_file.name.split(".")[-1].lower()

    # TXT
    if file_type == "txt":
        return uploaded_file.read().decode("utf-8")

    # DOCX
    elif file_type == "docx":
        doc = docx.Document(uploaded_file)
        return "\n".join([p.text for p in doc.paragraphs])

    # PDF
    elif file_type == "pdf":
        text = ""

        # Normal text extraction
        reader = PyPDF2.PdfReader(uploaded_file)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"

        # OCR pages if little text
        if len(text.strip()) < 50:
            uploaded_file.seek(0)
            images = convert_from_bytes(uploaded_file.read())

            for img in images:
                text += pytesseract.image_to_string(img)

        return text

    # PPTX
    elif file_type == "pptx":
        prs = Presentation(uploaded_file)
        text = ""

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

        return text

    # Images
    elif file_type in ["png", "jpg", "jpeg"]:
        image = Image.open(uploaded_file)
        return pytesseract.image_to_string(image)

    return ""